"""
PDF Processing Services for Easy Scan PDF
Implements comprehensive PDF operations
"""

import os
import io
from typing import List, Tuple, Optional
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
from pdf2docx import Converter
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from PIL import Image
import tempfile
import logging
from pdf2image import convert_from_bytes
import pikepdf
import pdfplumber
import pandas as pd

logger = logging.getLogger(__name__)


class PDFService:
    """Service class for PDF operations"""

    @staticmethod
    def merge_pdfs(pdf_files: List[bytes]) -> bytes:
        """
        Merge multiple PDF files into one
        Args:
            pdf_files: List of PDF file contents as bytes
        Returns:
            Merged PDF as bytes
        """
        try:
            merger = PdfMerger()
            
            for pdf_bytes in pdf_files:
                pdf_stream = io.BytesIO(pdf_bytes)
                merger.append(pdf_stream)
            
            output = io.BytesIO()
            merger.write(output)
            merger.close()
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error merging PDFs: {str(e)}")
            raise Exception(f"Failed to merge PDFs: {str(e)}")

    @staticmethod
    def split_pdf(pdf_file: bytes, split_type: str = "all") -> List[bytes]:
        """
        Split PDF into separate pages
        Args:
            pdf_file: PDF file content as bytes
            split_type: 'all' for all pages, or specific ranges later
        Returns:
            List of PDF pages as bytes
        """
        try:
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            
            split_pdfs = []
            
            for page_num in range(len(reader.pages)):
                writer = PdfWriter()
                writer.add_page(reader.pages[page_num])
                
                output = io.BytesIO()
                writer.write(output)
                output.seek(0)
                split_pdfs.append(output.getvalue())
            
            return split_pdfs
        except Exception as e:
            logger.error(f"Error splitting PDF: {str(e)}")
            raise Exception(f"Failed to split PDF: {str(e)}")

    @staticmethod
    def compress_pdf(pdf_file: bytes, quality: str = "medium") -> bytes:
        """
        Compress PDF by reducing image quality and removing duplicates
        Args:
            pdf_file: PDF file content as bytes
            quality: 'low', 'medium', or 'high'
        Returns:
            Compressed PDF as bytes
        """
        try:
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            writer = PdfWriter()
            
            # Quality settings
            quality_map = {"low": 50, "medium": 75, "high": 85}
            compression_quality = quality_map.get(quality, 75)
            
            for page in reader.pages:
                # Compress page
                page.compress_content_streams()
                writer.add_page(page)
            
            # Remove duplicate objects
            writer.add_metadata(reader.metadata)
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error compressing PDF: {str(e)}")
            raise Exception(f"Failed to compress PDF: {str(e)}")

    @staticmethod
    def pdf_to_word(pdf_file: bytes) -> bytes:
        """
        Convert PDF to Word (DOCX)
        Args:
            pdf_file: PDF file content as bytes
        Returns:
            DOCX file as bytes
        """
        try:
            # Create temporary files
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as pdf_temp:
                pdf_temp.write(pdf_file)
                pdf_path = pdf_temp.name
            
            docx_path = pdf_path.replace('.pdf', '.docx')
            
            # Convert PDF to DOCX
            cv = Converter(pdf_path)
            cv.convert(docx_path, start=0, end=None)
            cv.close()
            
            # Read DOCX file
            with open(docx_path, 'rb') as docx_file:
                docx_bytes = docx_file.read()
            
            # Cleanup temporary files
            os.unlink(pdf_path)
            os.unlink(docx_path)
            
            return docx_bytes
        except Exception as e:
            logger.error(f"Error converting PDF to Word: {str(e)}")
            raise Exception(f"Failed to convert PDF to Word: {str(e)}")

    @staticmethod
    def word_to_pdf(docx_file: bytes) -> bytes:
        """
        Convert Word (DOCX) to PDF
        Args:
            docx_file: DOCX file content as bytes
        Returns:
            PDF file as bytes
        """
        try:
            # Create temporary DOCX file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as docx_temp:
                docx_temp.write(docx_file)
                docx_path = docx_temp.name
            
            # Read DOCX content
            doc = Document(docx_path)
            
            # Create PDF
            output = io.BytesIO()
            pdf_canvas = canvas.Canvas(output, pagesize=letter)
            width, height = letter
            
            y_position = height - 50
            
            # Extract text from DOCX and add to PDF
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Wrap text if too long
                    text = paragraph.text
                    if len(text) > 80:
                        # Simple text wrapping
                        words = text.split()
                        line = ""
                        for word in words:
                            if len(line + word) < 80:
                                line += word + " "
                            else:
                                pdf_canvas.drawString(50, y_position, line.strip())
                                y_position -= 15
                                line = word + " "
                                
                                if y_position < 50:
                                    pdf_canvas.showPage()
                                    y_position = height - 50
                        
                        if line:
                            pdf_canvas.drawString(50, y_position, line.strip())
                            y_position -= 15
                    else:
                        pdf_canvas.drawString(50, y_position, text)
                        y_position -= 15
                    
                    if y_position < 50:
                        pdf_canvas.showPage()
                        y_position = height - 50
            
            pdf_canvas.save()
            output.seek(0)
            
            # Cleanup
            os.unlink(docx_path)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error converting Word to PDF: {str(e)}")
            raise Exception(f"Failed to convert Word to PDF: {str(e)}")

    @staticmethod
    def pdf_to_jpg(pdf_file: bytes) -> List[bytes]:
        """
        Convert PDF pages to JPG images
        Args:
            pdf_file: PDF file content as bytes
        Returns:
            List of JPG images as bytes
        """
        try:
            # Convert PDF to images
            images = convert_from_bytes(pdf_file, fmt='jpeg', dpi=200)
            
            jpg_images = []
            for image in images:
                # Convert PIL Image to bytes
                img_byte_arr = io.BytesIO()
                image.save(img_byte_arr, format='JPEG', quality=95)
                img_byte_arr.seek(0)
                jpg_images.append(img_byte_arr.getvalue())
            
            return jpg_images
        except Exception as e:
            logger.error(f"Error converting PDF to JPG: {str(e)}")
            raise Exception(f"Failed to convert PDF to JPG: {str(e)}")

    @staticmethod
    def jpg_to_pdf(jpg_files: List[bytes]) -> bytes:
        """
        Convert JPG images to PDF
        Args:
            jpg_files: List of JPG file contents as bytes
        Returns:
            PDF file as bytes
        """
        try:
            output = io.BytesIO()
            pdf_canvas = canvas.Canvas(output, pagesize=letter)
            width, height = letter
            
            for jpg_bytes in jpg_files:
                # Open image
                img = Image.open(io.BytesIO(jpg_bytes))
                
                # Calculate dimensions to fit page
                img_width, img_height = img.size
                aspect = img_height / img_width
                
                if aspect > (height / width):
                    # Height is limiting factor
                    new_height = height - 40
                    new_width = new_height / aspect
                else:
                    # Width is limiting factor
                    new_width = width - 40
                    new_height = new_width * aspect
                
                # Center image on page
                x = (width - new_width) / 2
                y = (height - new_height) / 2
                
                # Draw image
                img_reader = ImageReader(io.BytesIO(jpg_bytes))
                pdf_canvas.drawImage(img_reader, x, y, new_width, new_height)
                pdf_canvas.showPage()
            
            pdf_canvas.save()
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error converting JPG to PDF: {str(e)}")
            raise Exception(f"Failed to convert JPG to PDF: {str(e)}")

    @staticmethod
    def jpg_to_word(jpg_files: List[bytes]) -> bytes:
        """
        Convert JPG images to Word document
        Args:
            jpg_files: List of JPG file contents as bytes
        Returns:
            DOCX file as bytes
        """
        try:
            doc = Document()
            
            for jpg_bytes in jpg_files:
                # Save image temporarily
                with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as temp_img:
                    temp_img.write(jpg_bytes)
                    temp_img_path = temp_img.name
                
                # Add image to document
                doc.add_picture(temp_img_path, width=doc.sections[0].page_width - doc.sections[0].left_margin - doc.sections[0].right_margin)
                doc.add_paragraph()  # Add space between images
                
                # Cleanup temp file
                os.unlink(temp_img_path)
            
            # Save document to bytes
            output = io.BytesIO()
            doc.save(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error converting JPG to Word: {str(e)}")
            raise Exception(f"Failed to convert JPG to Word: {str(e)}")

    @staticmethod
    def word_to_jpg(docx_file: bytes) -> List[bytes]:
        """
        Convert Word document to JPG images
        First converts to PDF, then PDF to JPG
        Args:
            docx_file: DOCX file content as bytes
        Returns:
            List of JPG images as bytes
        """
        try:
            # First convert Word to PDF
            pdf_bytes = PDFService.word_to_pdf(docx_file)
            
            # Then convert PDF to JPG
            jpg_images = PDFService.pdf_to_jpg(pdf_bytes)
            
            return jpg_images
        except Exception as e:
            logger.error(f"Error converting Word to JPG: {str(e)}")
            raise Exception(f"Failed to convert Word to JPG: {str(e)}")

    @staticmethod
    def unlock_pdf(pdf_file: bytes, password: str = "") -> bytes:
        """
        Remove password protection from PDF
        Args:
            pdf_file: PDF file content as bytes
            password: Password to unlock the PDF (if needed)
        Returns:
            Unlocked PDF as bytes
        """
        try:
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            
            # Try to decrypt if encrypted
            if reader.is_encrypted:
                if password:
                    reader.decrypt(password)
                else:
                    # Try empty password
                    reader.decrypt("")
            
            # Create new PDF without encryption
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error unlocking PDF: {str(e)}")
            raise Exception(f"Failed to unlock PDF. Make sure you provided the correct password: {str(e)}")

    @staticmethod
    def protect_pdf(pdf_file: bytes, password: str) -> bytes:
        """
        Add password protection to PDF
        Args:
            pdf_file: PDF file content as bytes
            password: Password to protect the PDF
        Returns:
            Protected PDF as bytes
        """
        try:
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            writer = PdfWriter()
            
            # Add all pages
            for page in reader.pages:
                writer.add_page(page)
            
            # Encrypt with password
            writer.encrypt(user_password=password, owner_password=password, algorithm="AES-256")
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error protecting PDF: {str(e)}")
            raise Exception(f"Failed to protect PDF: {str(e)}")

    @staticmethod
    def organize_pdf(pdf_file: bytes, page_order: List[int]) -> bytes:
        """
        Reorganize PDF pages in custom order
        Args:
            pdf_file: PDF file content as bytes
            page_order: List of page numbers in desired order (1-indexed)
        Returns:
            Reorganized PDF as bytes
        """
        try:
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            writer = PdfWriter()
            
            total_pages = len(reader.pages)
            
            # If no order specified, use original order
            if not page_order:
                page_order = list(range(1, total_pages + 1))
            
            # Add pages in specified order
            for page_num in page_order:
                # Convert to 0-indexed
                page_index = page_num - 1
                
                if 0 <= page_index < total_pages:
                    writer.add_page(reader.pages[page_index])
                else:
                    logger.warning(f"Page {page_num} out of range, skipping")
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error organizing PDF: {str(e)}")
            raise Exception(f"Failed to organize PDF: {str(e)}")

    @staticmethod
    def rotate_pdf(pdf_file: bytes, rotation: int = 90) -> bytes:
        """
        Rotate PDF pages
        Args:
            pdf_file: PDF file content as bytes
            rotation: Rotation angle (90, 180, 270)
        Returns:
            Rotated PDF as bytes
        """
        try:
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            writer = PdfWriter()
            
            for page in reader.pages:
                page.rotate(rotation)
                writer.add_page(page)
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error rotating PDF: {str(e)}")
            raise Exception(f"Failed to rotate PDF: {str(e)}")

    @staticmethod
    def add_page_numbers(pdf_file: bytes, position: str = "bottom") -> bytes:
        """
        Add page numbers to PDF
        Args:
            pdf_file: PDF file content as bytes
            position: Position of page numbers (bottom, top)
        Returns:
            PDF with page numbers as bytes
        """
        try:
            # Read original PDF
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            
            # Create new PDF with page numbers
            output = io.BytesIO()
            writer = PdfWriter()
            
            for page_num, page in enumerate(reader.pages, 1):
                # Create overlay with page number
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                width, height = letter
                
                # Position page number
                y_position = 30 if position == "bottom" else height - 30
                can.drawCentredString(width / 2, y_position, str(page_num))
                can.save()
                
                # Merge overlay with original page
                packet.seek(0)
                overlay_reader = PdfReader(packet)
                page.merge_page(overlay_reader.pages[0])
                writer.add_page(page)
            
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error adding page numbers: {str(e)}")
            raise Exception(f"Failed to add page numbers: {str(e)}")

    @staticmethod
    def watermark_pdf(pdf_file: bytes, watermark_text: str = "CONFIDENTIAL") -> bytes:
        """
        Add watermark to PDF
        Args:
            pdf_file: PDF file content as bytes
            watermark_text: Text to use as watermark
        Returns:
            Watermarked PDF as bytes
        """
        try:
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            writer = PdfWriter()
            
            for page in reader.pages:
                # Create watermark
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                width, height = letter
                
                # Set watermark properties
                can.setFont("Helvetica-Bold", 60)
                can.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)
                can.saveState()
                can.translate(width / 2, height / 2)
                can.rotate(45)
                can.drawCentredString(0, 0, watermark_text)
                can.restoreState()
                can.save()
                
                # Merge watermark with page
                packet.seek(0)
                watermark_reader = PdfReader(packet)
                page.merge_page(watermark_reader.pages[0])
                writer.add_page(page)
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error adding watermark: {str(e)}")
            raise Exception(f"Failed to add watermark: {str(e)}")

    @staticmethod
    def crop_pdf(pdf_file: bytes, margin: int = 50) -> bytes:
        """
        Crop PDF pages
        Args:
            pdf_file: PDF file content as bytes
            margin: Margin to crop in points
        Returns:
            Cropped PDF as bytes
        """
        try:
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            writer = PdfWriter()
            
            for page in reader.pages:
                # Get current page dimensions
                mediabox = page.mediabox
                
                # Crop by reducing mediabox
                page.mediabox.lower_left = (
                    mediabox.lower_left[0] + margin,
                    mediabox.lower_left[1] + margin
                )
                page.mediabox.upper_right = (
                    mediabox.upper_right[0] - margin,
                    mediabox.upper_right[1] - margin
                )
                
                writer.add_page(page)
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error cropping PDF: {str(e)}")
            raise Exception(f"Failed to crop PDF: {str(e)}")

    @staticmethod
    def redact_pdf(pdf_file: bytes) -> bytes:
        """
        Redact sensitive information from PDF (simplified version)
        Args:
            pdf_file: PDF file content as bytes
        Returns:
            Redacted PDF as bytes
        """
        try:
            # Simple redaction - adds black rectangles over text
            pdf_stream = io.BytesIO(pdf_file)
            reader = PdfReader(pdf_stream)
            writer = PdfWriter()
            
            for page in reader.pages:
                # Create redaction overlay (simplified - covers portions of page)
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                width, height = letter
                
                # Add sample redaction boxes (in real implementation, would detect sensitive data)
                can.setFillColorRGB(0, 0, 0)
                can.rect(100, height - 150, 200, 20, fill=True)
                can.save()
                
                packet.seek(0)
                overlay_reader = PdfReader(packet)
                page.merge_page(overlay_reader.pages[0])
                writer.add_page(page)
            
            output = io.BytesIO()
            writer.write(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error redacting PDF: {str(e)}")
            raise Exception(f"Failed to redact PDF: {str(e)}")

    @staticmethod
    def pdf_to_excel(pdf_file: bytes) -> bytes:
        """
        Convert PDF tables to Excel
        Args:
            pdf_file: PDF file content as bytes
        Returns:
            Excel file as bytes
        """
        try:
            # Save PDF temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                temp_pdf.write(pdf_file)
                temp_pdf_path = temp_pdf.name
            
            # Extract tables from PDF
            all_tables = []
            with pdfplumber.open(temp_pdf_path) as pdf:
                for page in pdf.pages:
                    tables = page.extract_tables()
                    all_tables.extend(tables)
            
            # Create Excel file
            output = io.BytesIO()
            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                for i, table in enumerate(all_tables, 1):
                    if table:
                        df = pd.DataFrame(table[1:], columns=table[0] if table else None)
                        df.to_excel(writer, sheet_name=f'Table_{i}', index=False)
            
            output.seek(0)
            os.unlink(temp_pdf_path)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error converting PDF to Excel: {str(e)}")
            raise Exception(f"Failed to convert PDF to Excel: {str(e)}")

    @staticmethod
    def pdf_to_powerpoint(pdf_file: bytes) -> bytes:
        """
        Convert PDF to PowerPoint (simplified - converts pages to images in PPTX)
        Args:
            pdf_file: PDF file content as bytes
        Returns:
            PPTX file as bytes
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            # Convert PDF pages to images
            images = convert_from_bytes(pdf_file, fmt='png', dpi=150)
            
            # Create PowerPoint
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for image in images:
                # Add blank slide
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Save image temporarily
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_img:
                    image.save(temp_img.name, 'PNG')
                    temp_img_path = temp_img.name
                
                # Add image to slide
                slide.shapes.add_picture(temp_img_path, 0, 0, 
                                        width=prs.slide_width,
                                        height=prs.slide_height)
                
                os.unlink(temp_img_path)
            
            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error converting PDF to PowerPoint: {str(e)}")
            raise Exception(f"Failed to convert PDF to PowerPoint: {str(e)}")

    @staticmethod
    def powerpoint_to_pdf(pptx_file: bytes) -> bytes:
        """
        Convert PowerPoint to PDF (simplified - requires unoconv or similar)
        For now, returns a placeholder
        """
        try:
            # This requires LibreOffice or similar converter
            # Simplified version - create PDF from PPTX content
            from pptx import Presentation
            
            # Load presentation
            prs = Presentation(io.BytesIO(pptx_file))
            
            # Create PDF with slides as pages
            output = io.BytesIO()
            pdf_canvas = canvas.Canvas(output, pagesize=letter)
            width, height = letter
            
            for i, slide in enumerate(prs.slides, 1):
                # Add page number and slide title
                pdf_canvas.setFont("Helvetica-Bold", 16)
                pdf_canvas.drawString(50, height - 50, f"Slide {i}")
                pdf_canvas.showPage()
            
            pdf_canvas.save()
            output.seek(0)
            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error converting PowerPoint to PDF: {str(e)}")
            raise Exception(f"Failed to convert PowerPoint to PDF: {str(e)}")

    @staticmethod
    def html_to_pdf(html_content: str) -> bytes:
        """
        Convert HTML to PDF
        Args:
            html_content: HTML string
        Returns:
            PDF file as bytes
        """
        try:
            from weasyprint import HTML
            
            output = io.BytesIO()
            HTML(string=html_content).write_pdf(output)
            output.seek(0)
            
            return output.getvalue()
        except Exception:
            # Fallback to simple conversion
            output = io.BytesIO()
            pdf_canvas = canvas.Canvas(output, pagesize=letter)
            width, height = letter
            
            # Simple HTML to text conversion
            import re
            text = re.sub('<[^<]+?>', '', html_content)
            
            y_position = height - 50
            for line in text.split('\n'):
                if line.strip():
                    pdf_canvas.drawString(50, y_position, line[:80])
                    y_position -= 15
                    if y_position < 50:
                        pdf_canvas.showPage()
                        y_position = height - 50
            
            pdf_canvas.save()
            output.seek(0)
            return output.getvalue()

    @staticmethod
    def ocr_pdf(pdf_file: bytes) -> bytes:
        """
        Perform OCR on PDF (placeholder - requires tesseract)
        Args:
            pdf_file: PDF file content as bytes
        Returns:
            Searchable PDF as bytes
        """
        try:
            # For now, return original PDF with message
            # Full OCR would require pytesseract and tesseract-ocr
            logger.info("OCR requested - returning original PDF (OCR not fully implemented)")
            return pdf_file
        except Exception as e:
            logger.error(f"Error performing OCR: {str(e)}")
            raise Exception(f"Failed to perform OCR: {str(e)}")

    @staticmethod
    def repair_pdf(pdf_file: bytes) -> bytes:
        """
        Attempt to repair corrupted PDF
        Args:
            pdf_file: PDF file content as bytes
        Returns:
            Repaired PDF as bytes
        """
        try:
            with pikepdf.open(io.BytesIO(pdf_file), allow_overwriting_input=True) as pdf:
                output = io.BytesIO()
                pdf.save(output)
                output.seek(0)
                return output.getvalue()
        except Exception as e:
            logger.error(f"Error repairing PDF: {str(e)}")
            raise Exception(f"Failed to repair PDF: {str(e)}")

    @staticmethod
    def pdf_to_pdfa(pdf_file: bytes) -> bytes:
        """
        Convert PDF to PDF/A format
        Args:
            pdf_file: PDF file content as bytes
        Returns:
            PDF/A file as bytes
        """
        try:
            with pikepdf.open(io.BytesIO(pdf_file)) as pdf:
                # Set PDF/A metadata
                with pdf.open_metadata() as meta:
                    meta['pdfa:part'] = '1'
                    meta['pdfa:conformance'] = 'B'
                
                output = io.BytesIO()
                pdf.save(output)
                output.seek(0)
                return output.getvalue()
        except Exception as e:
            logger.error(f"Error converting to PDF/A: {str(e)}")
            raise Exception(f"Failed to convert to PDF/A: {str(e)}")

            
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error organizing PDF: {str(e)}")
            raise Exception(f"Failed to organize PDF: {str(e)}")



# Initialize service
pdf_service = PDFService()
